"""
Genera el PPT semanal de INGRESOS (Activación) de una campaña.

Estructura esperada del template "Template_Reporte_Ingresos":
  Slide 1 → portada
  Slide 2 → plantilla de detalle (se duplica una vez por fila)

Se reutilizan los helpers de generar_ppt.py (duplicación de slides con
relationships y fondo, normalización/compresión de imágenes, cálculo de
posiciones de fotos). Ese archivo NO se modifica: acá solo se importa.

PLACEHOLDERS SOPORTADOS
-----------------------
Se reemplazan todos los que existan en el template; los que no estén,
simplemente se ignoran. Así podés ajustar el template sin tocar el código.

Portada:
  [CAMPAÑA] [CAMPANA] [CLIENTE] [SEMANA] [N SEMANA] [AÑO] [ANIO]
  [PERIODO] [FECHA ENVIO DEL REPORTE] [FECHA] [TOTAL SALAS] [TOTAL REGISTROS]

Detalle (por fila):
  [CÓD.] - [NOMBRE SALA]              (combinado, igual que en Implementación)
  [DIRECCIÓN], [COMUNA] - [REGION]    (combinado)
  [CÓD.] [COD] [NOMBRE SALA] [DIRECCIÓN] [COMUNA] [REGION] [CADENA]
  [MARCA] [CATEGORIA] [SKU] [MATERIAL] [CANTIDAD] [GUIA]
  [PROCESO] [DETALLE] [OBSERVACIONES] [EJECUTIVO]
  [ACTIVIDAD] [TIPO ACTIVIDAD] [SEMANA]
  [FECHA INICIO] [FECHA TERMINO] [FECHA ENTREGA] [FECHA COMPROMISO]
  [HORA INICIO] [HORA TERMINO]
  [FOTO 1] .. [FOTO 4]
"""
import os
from io import BytesIO

from pptx import Presentation
from PIL import Image
from googleapiclient.http import MediaIoBaseDownload

from config import IDX
from config_ingresos import (
    NOMBRE_TEMPLATE_INGRESOS,
    NOMBRE_CARPETA_TEMPLATES,
    TEMPLATE_PPT_INGRESOS_ID,
    MAX_FOTOS_INGRESOS,
    FOTOS_POR_PROCESO_INGRESOS,
)
from utils import fmt, buscar_foto_blob
from google_clients import get_drive

# Helpers reutilizados del flujo de Implementación (no se duplica lógica)
from generar_ppt import (
    _normalizar_imagen,
    _reemplazar_texto,
    _encontrar_placeholders_fotos,
    _calcular_posiciones,
    _duplicar_slide,
    _limpiar_placeholder_foto,
)

MIME_SLIDES = "application/vnd.google-apps.presentation"
MIME_PPTX = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
MIME_FOLDER = "application/vnd.google-apps.folder"

# Cache: el template es el mismo para todas las campañas de la corrida.
_template_cache = None


# ============================================================
# BÚSQUEDA Y DESCARGA DEL TEMPLATE
# ============================================================
def _buscar_template_ingresos(drive):
    """
    Ubica el archivo del template. Orden de intentos:
      1. TEMPLATE_PPT_INGRESOS_ID en config_ingresos.py
      2. Variable de entorno INGRESOS_TEMPLATE_ID
      3. Dentro de la carpeta "Templates_Reportes" (Mi unidad)
      4. Búsqueda global por nombre

    Devuelve (file_id, mime_type).
    """
    id_fijo = (TEMPLATE_PPT_INGRESOS_ID or os.environ.get("INGRESOS_TEMPLATE_ID") or "").strip()
    if id_fijo:
        info = drive.files().get(
            fileId=id_fijo, fields="id, name, mimeType", supportsAllDrives=True
        ).execute()
        print(f"  [TEMPLATE] ID fijo: '{info['name']}' ({info['id']})")
        return info["id"], info["mimeType"]

    nombre_q = NOMBRE_TEMPLATE_INGRESOS.replace("'", "\\'")
    print(f"  [TEMPLATE] Buscando '{NOMBRE_TEMPLATE_INGRESOS}'...")

    # --- Intento 1: dentro de la carpeta Templates_Reportes ---
    try:
        carpeta_q = NOMBRE_CARPETA_TEMPLATES.replace("'", "\\'")
        resp = drive.files().list(
            q=f"name = '{carpeta_q}' and mimeType = '{MIME_FOLDER}' and trashed = false",
            fields="files(id, name)",
            supportsAllDrives=True, includeItemsFromAllDrives=True,
        ).execute()
        for carpeta in resp.get("files", []):
            resp2 = drive.files().list(
                q=(f"name = '{nombre_q}' and '{carpeta['id']}' in parents "
                   f"and trashed = false"),
                fields="files(id, name, mimeType)",
                supportsAllDrives=True, includeItemsFromAllDrives=True,
            ).execute()
            archivos = resp2.get("files", [])
            if archivos:
                a = archivos[0]
                print(f"  [TEMPLATE] ✓ En '{NOMBRE_CARPETA_TEMPLATES}': {a['name']} ({a['id']})")
                return a["id"], a["mimeType"]
        print(f"  [TEMPLATE] No está en '{NOMBRE_CARPETA_TEMPLATES}'. Busco global...")
    except Exception as e:
        print(f"  [TEMPLATE] Error buscando en la carpeta: {e}. Busco global...")

    # --- Intento 2: búsqueda global por nombre ---
    resp3 = drive.files().list(
        q=f"name = '{nombre_q}' and trashed = false",
        fields="files(id, name, mimeType)",
        supportsAllDrives=True, includeItemsFromAllDrives=True,
    ).execute()
    archivos3 = resp3.get("files", [])
    if archivos3:
        a = archivos3[0]
        print(f"  [TEMPLATE] ✓ Encontrado (global): {a['name']} ({a['id']})")
        return a["id"], a["mimeType"]

    raise RuntimeError(
        f"No se encontró '{NOMBRE_TEMPLATE_INGRESOS}'. Revisá que la carpeta "
        f"'{NOMBRE_CARPETA_TEMPLATES}' (o el archivo) esté compartida con la "
        f"cuenta de servicio, o pegá el ID en TEMPLATE_PPT_INGRESOS_ID."
    )


def _descargar_template(drive, file_id, mime):
    """Baja el template como .pptx en memoria (exporta si es Google Slides)."""
    if mime == MIME_SLIDES:
        req = drive.files().export_media(fileId=file_id, mimeType=MIME_PPTX)
    else:
        req = drive.files().get_media(fileId=file_id, supportsAllDrives=True)
    buf = BytesIO()
    downloader = MediaIoBaseDownload(buf, req)
    done = False
    while not done:
        _, done = downloader.next_chunk()
    buf.seek(0)
    return buf


def _template_bytes(drive):
    """Devuelve los bytes del template, cacheados para toda la corrida."""
    global _template_cache
    if _template_cache is None:
        file_id, mime = _buscar_template_ingresos(drive)
        _template_cache = _descargar_template(drive, file_id, mime).getvalue()
        print(f"  [TEMPLATE] Descargado ({len(_template_cache)/1024:.0f} KB)")
    return BytesIO(_template_cache)


# ============================================================
# TOKENS
# ============================================================
def _tokens_portada(campana, cliente, etiqueta_semana, semana, anio,
                    lunes, domingo, hoy, filas):
    salas = {str(r[IDX["COD"]] or "").strip() for r in filas}
    salas.discard("")
    periodo = f"{lunes.strftime('%d/%m/%Y')} al {domingo.strftime('%d/%m/%Y')}"
    return [
        ("[FECHA ENVIO DEL REPORTE]", hoy.strftime("%d/%m/%Y")),
        ("[TOTAL REGISTROS]", str(len(filas))),
        ("[TOTAL SALAS]", str(len(salas))),
        ("[N SEMANA]", str(semana)),
        ("[CAMPAÑA]", campana),
        ("[CAMPANA]", campana),
        ("[CLIENTE]", cliente),
        ("[PERIODO]", periodo),
        ("[SEMANA]", etiqueta_semana),
        ("[AÑO]", str(anio)),
        ("[ANIO]", str(anio)),
        ("[FECHA]", hoy.strftime("%d/%m/%Y")),
    ]


def _tokens_fila(row, etiqueta_semana, proceso):
    """
    Tokens de una slide de detalle.

    OJO CON EL ORDEN: los combinados van PRIMERO. Si reemplazáramos [CÓD.]
    antes que "[CÓD.] - [NOMBRE SALA]", el combinado ya no existiría y la
    línea quedaría a medias.
    """
    g = lambda k: str(row[IDX[k]] or "")  # noqa: E731
    cod, sala = g("COD"), g("NOMBRE_SALA")
    direccion, comuna, region = g("DIRECCION"), g("COMUNA"), g("REGION")

    return [
        # --- combinados (primero) ---
        ("[CÓD.] - [NOMBRE SALA]", f"{cod} - {sala}"),
        ("[COD] - [NOMBRE SALA]", f"{cod} - {sala}"),
        ("[DIRECCIÓN], [COMUNA] - [REGION]", f"{direccion}, {comuna} - {region}"),
        ("[DIRECCION], [COMUNA] - [REGION]", f"{direccion}, {comuna} - {region}"),
        # --- individuales ---
        ("[TIPO ACTIVIDAD]", g("TIPO_ACTIVIDAD")),
        ("[NOMBRE SALA]", sala),
        ("[FECHA COMPROMISO]", fmt(row[IDX["FECHA_COMPROMISO"]])),
        ("[FECHA TERMINO]", fmt(row[IDX["FECHA_TERMINO"]])),
        ("[FECHA ENTREGA]", fmt(row[IDX["FECHA_ENTREGA"]])),
        ("[FECHA INICIO]", fmt(row[IDX["FECHA_INICIO"]])),
        ("[HORA TERMINO]", g("HORA_TERMINO")),
        ("[HORA INICIO]", g("HORA_INICIO")),
        ("[OBSERVACIONES]", g("OBSERVACIONES")),
        ("[DIRECCIÓN]", direccion),
        ("[DIRECCION]", direccion),
        ("[CATEGORIA]", g("CATEGORIA")),
        ("[CATEGORÍA]", g("CATEGORIA")),
        ("[EJECUTIVO]", g("EJECUTIVO")),
        ("[ACTIVIDAD]", g("ACTIVIDAD")),
        ("[CANTIDAD]", g("CANTIDAD")),
        ("[MATERIAL]", g("MATERIAL")),
        ("[PROCESO]", proceso),
        ("[DETALLE]", g("DETALLE")),
        ("[CADENA]", g("CADENA")),
        ("[COMUNA]", comuna),
        ("[REGION]", region),
        ("[REGIÓN]", region),
        ("[SEMANA]", etiqueta_semana),
        ("[MARCA]", g("MARCA")),
        ("[GUIA]", g("GUIA")),
        ("[GUÍA]", g("GUIA")),
        ("[CÓD.]", cod),
        ("[COD]", cod),
        ("[SKU]", g("SKU")),
    ]


def _orden_codigo(row):
    """Orden numérico real de la sala: 'J9' antes que 'J633'."""
    import re
    cod = str(row[IDX["COD"]] or "").strip()
    m = re.match(r"^([A-Za-z]*)(\d+)?", cod)
    if m:
        return (m.group(1) or "", int(m.group(2)) if m.group(2) else 0)
    return (cod, 0)


def _proceso_limpio(row):
    """Normaliza 'Reagenda interna/externa' → 'Reagenda' (igual que Implementación)."""
    p = str(row[IDX["PROCESO"]] or "").strip()
    return "Reagenda" if p.lower().startswith("reagenda") else p


def _rutas_fotos(row, proceso):
    """Rutas de foto de la fila, sin repetidas, respetando el máximo."""
    tope = FOTOS_POR_PROCESO_INGRESOS.get(proceso.strip().lower(), MAX_FOTOS_INGRESOS)
    rutas = []
    for k in ("FOTO1", "FOTO2", "FOTO3", "FOTO4"):
        ruta = (row[IDX[k]] or "").strip()
        if ruta and ruta not in rutas:
            rutas.append(ruta)
    return rutas[:min(tope, MAX_FOTOS_INGRESOS)]


# ============================================================
# GENERACIÓN
# ============================================================
def generar_ppt_ingresos(campana, filas, etiqueta_semana, semana, anio,
                         lunes, domingo, hoy):
    """Genera el PPT semanal de ingresos de una campaña. Devuelve bytes."""
    drive = get_drive()
    cliente = str(filas[0][IDX["CLIENTE"]] or "")
    pres = Presentation(_template_bytes(drive))

    if len(pres.slides) < 2:
        raise RuntimeError(
            f"El template '{NOMBRE_TEMPLATE_INGRESOS}' tiene "
            f"{len(pres.slides)} slide(s). Se necesitan 2: portada + detalle."
        )

    # ---- PORTADA ----
    portada = pres.slides[0]
    for buscar, reemplazar in _tokens_portada(
        campana, cliente, etiqueta_semana, semana, anio, lunes, domingo, hoy, filas
    ):
        _reemplazar_texto(portada, buscar, reemplazar)

    # ---- DETALLE: una slide por fila ----
    filas_ordenadas = sorted(filas, key=_orden_codigo)
    print(f"  Slides a generar: {len(filas_ordenadas)}")

    template_detalle = pres.slides[1]
    slides = [template_detalle]
    for _ in range(len(filas_ordenadas) - 1):
        slides.append(_duplicar_slide(pres, template_detalle))

    for i, row in enumerate(filas_ordenadas):
        slide = slides[i]
        proceso = _proceso_limpio(row)
        cod = str(row[IDX["COD"]] or "")
        sala = str(row[IDX["NOMBRE_SALA"]] or "")
        print(f"    Slide {i+1}/{len(filas_ordenadas)}: {cod} {sala} | {proceso}")

        for buscar, reemplazar in _tokens_fila(row, etiqueta_semana, proceso):
            _reemplazar_texto(slide, buscar, reemplazar)

        rutas = _rutas_fotos(row, proceso)

        placeholders = _encontrar_placeholders_fotos(slide)
        info = [(sh.left, sh.top, sh.width, sh.height)
                for _, sh in sorted(placeholders.items())]
        for sh in placeholders.values():
            _limpiar_placeholder_foto(sh)

        if not rutas:
            continue

        posiciones = _calcular_posiciones(len(rutas), info)
        for j, ruta in enumerate(rutas):
            if j >= len(posiciones):
                break
            area_left, area_top, area_w, area_h = posiciones[j]
            blob = buscar_foto_blob(drive, ruta)
            if blob is None:
                print(f"      [WARN] Foto {j+1} no encontrada: {ruta}")
                continue
            try:
                norm = _normalizar_imagen(blob)
                if norm is None:
                    print(f"      [SKIP] Foto {j+1} omitida")
                    continue
                norm.seek(0)
                img_w, img_h = Image.open(norm).size
                norm.seek(0)

                # "Fit" dentro del área, manteniendo proporción y centrado
                if (img_w / img_h) > (area_w / area_h):
                    final_w, final_h = area_w, int(area_w / (img_w / img_h))
                else:
                    final_h, final_w = area_h, int(area_h * (img_w / img_h))

                slide.shapes.add_picture(
                    norm,
                    area_left + (area_w - final_w) // 2,
                    area_top + (area_h - final_h) // 2,
                    width=final_w, height=final_h,
                )
                print(f"      OK foto {j+1}")
            except Exception as e:
                print(f"      [ERROR] foto {j+1}: {e}")

    buf = BytesIO()
    pres.save(buf)
    return buf.getvalue()
