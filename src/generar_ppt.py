"""
Genera el PPT del reporte de fotos.
Equivalente a generarSlides() en Apps Script, pero con python-pptx (mucho más rápido).

ESTRATEGIA:
1. Descargar el template Google Slides como .pptx (vía Drive API export).
2. Abrirlo con python-pptx.
3. La primera slide es la portada, la segunda es la plantilla de sala (se duplica N veces).
4. Reemplazar placeholders [CAMPAÑA], [FOTO 1], etc.
5. Insertar las imágenes en las posiciones donde estaban los placeholders.
"""
import copy
from io import BytesIO

from pptx import Presentation
from pptx.util import Emu
from PIL import Image, ImageOps

# Registrar soporte para HEIC (fotos de iPhone modernas)
try:
    from pillow_heif import register_heif_opener
    register_heif_opener()
except ImportError:
    pass  # Si no está instalado, seguimos sin soporte HEIC

from googleapiclient.http import MediaIoBaseDownload

from config import IDX, TEMPLATE_PPT_ID
from utils import fmt, san, buscar_foto_blob
from google_clients import get_drive

# Configuración de compresión de imágenes (más agresiva)
MAX_LADO_PX = 1400       # máximo en el lado más largo (antes 1600)
CALIDAD_JPEG = 75        # antes 82, ahora 75 (sigue viéndose perfecto en PPT)


def _normalizar_imagen(blob):
    """
    SIEMPRE convierte a JPEG, redimensiona y comprime.
    Soporta: JPEG, PNG, MPO, HEIC, TIFF, BMP, GIF, WebP.
    Devuelve un BytesIO listo para python-pptx.
    Si falla completamente, devuelve None (la foto se omite).
    """
    tam_original = 0
    if hasattr(blob, "getbuffer"):
        tam_original = blob.getbuffer().nbytes
    elif hasattr(blob, "tell") and hasattr(blob, "seek"):
        pos = blob.tell()
        blob.seek(0, 2)
        tam_original = blob.tell()
        blob.seek(pos)

    try:
        blob.seek(0)
        img = Image.open(blob)
        formato_original = getattr(img, "format", "?")

        # MPO viene multi-frame; nos quedamos con el primer frame
        if formato_original == "MPO":
            img.seek(0)

        # Cargar la imagen completa en memoria (necesario para algunos formatos)
        img.load()

        # Convertir a RGB (descarta alpha, CMYK, paleta, etc.)
        if img.mode != "RGB":
            img = img.convert("RGB")

        # Aplicar rotación según EXIF (fotos de celular suelen venir rotadas)
        try:
            img = ImageOps.exif_transpose(img)
        except Exception:
            pass

        # Redimensionar si excede el máximo
        w, h = img.size
        lado = max(w, h)
        if lado > MAX_LADO_PX:
            ratio = MAX_LADO_PX / lado
            img = img.resize((int(w * ratio), int(h * ratio)), Image.LANCZOS)

        # Guardar como JPEG optimizado SIEMPRE (no devolver el blob original)
        out = BytesIO()
        img.save(out, format="JPEG", quality=CALIDAD_JPEG, optimize=True, progressive=True)
        out.seek(0)

        tam_final = out.getbuffer().nbytes
        if tam_original > 0:
            ratio_red = tam_final / tam_original * 100
            print(f"        [{formato_original}] {tam_original/1024:.0f}KB → {tam_final/1024:.0f}KB ({ratio_red:.0f}%)")

        return out

    except Exception as e:
        print(f"      [ERROR-NORM] No se pudo normalizar imagen ({type(e).__name__}: {e}). Se omite.")
        return None


def _descargar_template_como_pptx():
    """Exporta el Google Slides template a .pptx en memoria."""
    drive = get_drive()
    req = drive.files().export_media(
        fileId=TEMPLATE_PPT_ID,
        mimeType="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )
    buf = BytesIO()
    downloader = MediaIoBaseDownload(buf, req)
    done = False
    while not done:
        _, done = downloader.next_chunk()
    buf.seek(0)
    return buf


def _iter_text_frames(slide):
    """Genera todos los text_frames de la slide (incluyendo dentro de tablas)."""
    for shape in slide.shapes:
        if shape.has_text_frame:
            yield shape, shape.text_frame
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    yield shape, cell.text_frame


def _reemplazar_texto(slide, buscar, reemplazar):
    """Reemplaza texto en toda la slide, preservando formato del primer run."""
    reemplazar = str(reemplazar or "")
    for _, tf in _iter_text_frames(slide):
        for para in tf.paragraphs:
            full = "".join(run.text for run in para.runs)
            if buscar in full:
                nuevo = full.replace(buscar, reemplazar)
                # Limpiar todos los runs menos el primero, conservando su formato
                if para.runs:
                    para.runs[0].text = nuevo
                    for run in para.runs[1:]:
                        run.text = ""


def _encontrar_placeholders_fotos(slide):
    """
    Busca las shapes que contienen [FOTO 1]..[FOTO 4] y devuelve sus posiciones.
    Retorna dict {1: shape, 2: shape, ...}
    """
    encontrados = {}
    for shape in list(slide.shapes):
        if not shape.has_text_frame:
            continue
        txt = shape.text_frame.text
        for i in range(1, 5):
            if f"[FOTO {i}]" in txt:
                encontrados[i] = shape
                break
    return encontrados


def _calcular_posiciones(n_fotos, placeholders_info):
    """
    Calcula posiciones (left, top, width, height) según cantidad real de fotos.
    Espejo de calcularPosicionesFotos() en Apps Script.
    placeholders_info: lista de (left, top, width, height) en EMU.
    """
    if n_fotos == 0 or not placeholders_info:
        return []

    lefts = [p[0] for p in placeholders_info]
    tops = [p[1] for p in placeholders_info]
    rights = [p[0] + p[2] for p in placeholders_info]
    bottoms = [p[1] + p[3] for p in placeholders_info]

    min_x, min_y = min(lefts), min(tops)
    max_x, max_y = max(rights), max(bottoms)
    total_w, total_h = max_x - min_x, max_y - min_y
    gap = Emu(80000)  # ~8pt de separación

    if n_fotos == 1:
        return [(min_x, min_y, total_w, total_h)]
    if n_fotos == 2:
        w = (total_w - gap) // 2
        return [
            (min_x, min_y, w, total_h),
            (min_x + w + gap, min_y, w, total_h),
        ]
    if n_fotos == 3:
        if len(placeholders_info) == 3:
            return placeholders_info[:3]
        w = (total_w - gap * 2) // 3
        return [(min_x + i * (w + gap), min_y, w, total_h) for i in range(3)]
    # 4 fotos: grilla 2x2
    w = (total_w - gap) // 2
    h = (total_h - gap) // 2
    return [
        (min_x, min_y, w, h),
        (min_x + w + gap, min_y, w, h),
        (min_x, min_y + h + gap, w, h),
        (min_x + w + gap, min_y + h + gap, w, h),
    ]


def _duplicar_slide(pres, slide_origen):
    """
    Duplica una slide. python-pptx no tiene .duplicate() nativo.
    
    IMPORTANTE: además de copiar las shapes (XML), hay que copiar las
    'relationships' de imágenes incrustadas (el logo, fotos, etc.).
    Si no, las imágenes aparecen como referencias rotas en la slide nueva.
    """
    blank_layout = slide_origen.slide_layout
    nueva = pres.slides.add_slide(blank_layout)
    # Borrar los placeholders del layout en la nueva slide
    for shape in list(nueva.shapes):
        sp = shape._element
        sp.getparent().remove(sp)

    # Construir mapa de rId del slide origen → rId del slide nuevo
    # para imágenes, links y otros recursos referenciados.
    rels_origen = slide_origen.part.rels
    rid_map = {}  # {rId_viejo: rId_nuevo}
    for rId, rel in rels_origen.items():
        if rel.is_external:
            # Hyperlinks externos (target_ref es la URL)
            nuevo_rId = nueva.part.relate_to(rel.target_ref, rel.reltype, is_external=True)
        else:
            # Imágenes y otros recursos internos: reusar el mismo target_part
            nuevo_rId = nueva.part.relate_to(rel.target_part, rel.reltype)
        rid_map[rId] = nuevo_rId

    # Copiar shapes del slide origen, traduciendo rIds en el XML
    from lxml import etree
    NS_R = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed"
    NS_R_LINK = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}link"
    NS_R_ID = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"

    for shape in slide_origen.shapes:
        el = copy.deepcopy(shape._element)
        # Recorrer todos los atributos 'r:embed', 'r:link', 'r:id' del subárbol
        # y reemplazar el rId viejo por el nuevo
        for descendant in el.iter():
            for attr_name in (NS_R, NS_R_LINK, NS_R_ID):
                if attr_name in descendant.attrib:
                    viejo = descendant.attrib[attr_name]
                    if viejo in rid_map:
                        descendant.attrib[attr_name] = rid_map[viejo]
        nueva.shapes._spTree.insert_element_before(el, "p:extLst")

    return nueva


def _mover_slide_a_posicion(pres, slide, posicion):
    """Mueve una slide a la posición dada (0-indexed)."""
    xml_slides = pres.slides._sldIdLst
    slides_list = list(xml_slides)
    # Encontrar el elemento de esta slide
    for el in slides_list:
        if el.get(
            "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
        ) == slide.part.partname:
            # Difícil identificar 1:1; usamos el rId
            pass
    # Más simple: trabajamos con el último (recién agregado) y lo movemos
    new_el = slides_list[-1]
    xml_slides.remove(new_el)
    xml_slides.insert(posicion, new_el)


def _limpiar_placeholder_foto(shape):
    """Borra el shape del placeholder de foto."""
    sp = shape._element
    sp.getparent().remove(sp)


def generar_ppt(campana, filas, hoy):
    """
    Genera el PPT y devuelve los bytes.
    Equivalente a generarSlides() en Apps Script.
    """
    drive = get_drive()
    template_buf = _descargar_template_como_pptx()
    pres = Presentation(template_buf)

    fecha_str = hoy.strftime("%d/%m/%Y")
    cliente = str(filas[0][IDX["CLIENTE"]] or "")

    # ---- PORTADA (slide 0) ----
    portada = pres.slides[0]
    _reemplazar_texto(portada, "[FECHA ENVIO DEL REPORTE]", fecha_str)
    _reemplazar_texto(portada, "[CAMPAÑA]", campana)
    _reemplazar_texto(portada, "[CLIENTE]", cliente)

    # ---- ORDENAR FILAS POR CÓDIGO NUMÉRICO DE SALA ----
    # Cada fila genera una slide (sala+material). Si la misma sala aparece
    # con varios materiales, son varias slides distintas.
    def _orden_codigo(row):
        cod = str(row[IDX["COD"]] or "").strip()
        # Intentar parsear como número para orden numérico real
        # Si tiene prefijo (ej. "J633"), separar parte alfa y parte numérica
        import re
        m = re.match(r"^([A-Za-z]*)(\d+)?", cod)
        if m:
            prefix = m.group(1) or ""
            numero = int(m.group(2)) if m.group(2) else 0
            return (prefix, numero)
        return (cod, 0)

    filas_ordenadas = sorted(filas, key=_orden_codigo)
    print(f"  Slides a generar: {len(filas_ordenadas)}")

    # Slide template de sala (la segunda en el template)
    template_sala = pres.slides[1]

    # Lista de slides ya creadas para sala (la primera reutiliza template_sala)
    slides_de_sala = [template_sala]

    # Duplicar para las filas adicionales
    for _ in range(len(filas_ordenadas) - 1):
        nueva = _duplicar_slide(pres, template_sala)
        slides_de_sala.append(nueva)

    # ---- LLENAR CADA SLIDE (1 slide por fila = sala+material) ----
    for fila_idx, row in enumerate(filas_ordenadas):
        slide = slides_de_sala[fila_idx]
        cod = str(row[IDX["COD"]] or "")
        nombre_sala = str(row[IDX["NOMBRE_SALA"]] or "")
        material = str(row[IDX["MATERIAL"]] or "")
        print(f"    Slide {fila_idx + 1}/{len(filas_ordenadas)}: {cod} {nombre_sala} | {material}")

        # Reemplazar textos
        _reemplazar_texto(
            slide, "[CÓD.] - [NOMBRE SALA]",
            f"{cod} - {nombre_sala}",
        )
        _reemplazar_texto(
            slide, "[DIRECCIÓN], [COMUNA] - [REGION]",
            f"{row[IDX['DIRECCION']] or ''}, {row[IDX['COMUNA']] or ''} - {row[IDX['REGION']] or ''}",
        )
        _reemplazar_texto(slide, "[MATERIAL]", material)
        _reemplazar_texto(slide, "[CANTIDAD]", str(row[IDX["CANTIDAD"]] or ""))
        # Normalizar Reagenda interna/externa → Reagenda
        proceso_raw = str(row[IDX["PROCESO"]] or "").strip()
        if proceso_raw.lower().startswith("reagenda"):
            proceso_raw = "Reagenda"
        _reemplazar_texto(slide, "[PROCESO]", proceso_raw)
        _reemplazar_texto(slide, "[FECHA ENTREGA]", fmt(row[IDX["FECHA_ENTREGA"]]))

        # Recolectar rutas de las 4 fotos de ESTA fila específica (max 4)
        rutas = []
        for fi in [IDX["FOTO1"], IDX["FOTO2"], IDX["FOTO3"], IDX["FOTO4"]]:
            ruta = (row[fi] or "").strip()
            if ruta and ruta not in rutas:
                rutas.append(ruta)
        rutas = rutas[:4]

        # Localizar placeholders de foto y sus posiciones
        placeholders = _encontrar_placeholders_fotos(slide)
        placeholders_info = []
        for i in sorted(placeholders.keys()):
            sh = placeholders[i]
            placeholders_info.append((sh.left, sh.top, sh.width, sh.height))

        # Borrar todos los placeholders de foto (los reemplazaremos por imágenes)
        for sh in placeholders.values():
            _limpiar_placeholder_foto(sh)

        # Si no hay fotos, la slide queda con los datos solamente (sin imágenes)
        if not rutas:
            continue

        # Calcular posiciones según cantidad real de fotos
        posiciones = _calcular_posiciones(len(rutas), placeholders_info)

        # Insertar fotos manteniendo su proporción dentro del área disponible
        for i, ruta in enumerate(rutas):
            if i >= len(posiciones):
                break
            area_left, area_top, area_w, area_h = posiciones[i]
            blob = buscar_foto_blob(drive, ruta)
            if blob is None:
                print(f"      [WARN] Foto {i+1} no encontrada: {ruta}")
                continue
            try:
                blob_normalizado = _normalizar_imagen(blob)
                if blob_normalizado is None:
                    print(f"      [SKIP] Foto {i+1} omitida (no se pudo procesar)")
                    continue

                # Obtener dimensiones reales de la imagen (en px) para calcular ratio
                blob_normalizado.seek(0)
                img_px = Image.open(blob_normalizado)
                img_w_px, img_h_px = img_px.size
                blob_normalizado.seek(0)  # rebobinar para que pptx pueda leerla

                # Calcular el tamaño final manteniendo proporción ("fit dentro del área")
                ratio_img = img_w_px / img_h_px
                ratio_area = area_w / area_h
                if ratio_img > ratio_area:
                    # Imagen más ancha que el área: ajustar por ancho
                    final_w = area_w
                    final_h = int(area_w / ratio_img)
                else:
                    # Imagen más alta que el área: ajustar por alto
                    final_h = area_h
                    final_w = int(area_h * ratio_img)

                # Centrar dentro del área
                final_left = area_left + (area_w - final_w) // 2
                final_top = area_top + (area_h - final_h) // 2

                slide.shapes.add_picture(
                    blob_normalizado,
                    final_left, final_top,
                    width=final_w, height=final_h,
                )
                print(f"      OK foto {i+1}")
            except Exception as e:
                print(f"      [ERROR] foto {i+1}: {e}")

    # Guardar a bytes
    buf = BytesIO()
    pres.save(buf)
    buf.seek(0)
    return buf.getvalue()
