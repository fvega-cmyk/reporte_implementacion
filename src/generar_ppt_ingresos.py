"""
Genera el PPT semanal de INGRESOS (Activación) de una campaña.

DIFERENCIA CLAVE CON IMPLEMENTACIÓN
-----------------------------------
En Implementación cada fila (sala + material) genera su propia slide.
Acá NO: las fotos de una activación se repiten fila por fila, así que se
agrupa por LOCAL y cada local es UNA slide:

    1 LOCAL = 1 SLIDE
      ├── cabecera del local (cód., sala, dirección, comuna, región...)
      ├── listado de materiales: [MATERIAL] ([CANTIDAD]) : [PROCESO]
      │     una línea por material del local
      └── 2 fotos, quitando duplicados entre todas las filas del local

La agrupación es por ID_LOCAL (config_ingresos.AGRUPAR_POR).

Estructura esperada del template "Template_Reporte_Ingresos":
  Slide 1 → portada
  Slide 2 → plantilla de local (se duplica una vez por local)

Se reutilizan los helpers de generar_ppt.py (duplicación de slides con
relationships y fondo, normalización/compresión de imágenes, cálculo de
posiciones de fotos). Ese archivo NO se modifica: acá solo se importa.

PLACEHOLDERS SOPORTADOS
-----------------------
Se reemplazan los que existan; los que no estén se ignoran.

Portada:
  [CAMPAÑA] [CAMPANA] [CLIENTE] [SEMANA] [N SEMANA] [AÑO] [ANIO] [PERIODO]
  [FECHA ENVIO DEL REPORTE] [FECHA] [TOTAL LOCALES] [TOTAL SALAS]
  [TOTAL MATERIALES] [TOTAL REGISTROS]

Slide de local:
  [CÓD.] - [NOMBRE SALA]              (combinado)
  [DIRECCIÓN], [COMUNA] - [REGION]    (combinado)
  [CÓD.] [COD] [ID LOCAL] [NOMBRE SALA] [DIRECCIÓN] [COMUNA] [REGION] [CADENA]
  [MARCA] [CATEGORIA] [EJECUTIVO] [ACTIVIDAD] [TIPO ACTIVIDAD] [SEMANA]
  [FECHA INICIO] [FECHA TERMINO] [FECHA ENTREGA] [FECHA COMPROMISO]
  [HORA INICIO] [HORA TERMINO] [OBSERVACIONES] [DETALLE]
  [FOTO 1] [FOTO 2]

  Listado de materiales (una línea por material):
  [MATERIAL] ([CANTIDAD]) : [PROCESO]
     ↑ poné esta línea en el template UNA vez, con el formato que quieras
       (viñeta, tamaño, color). El código la duplica por cada material,
       conservando ese formato, y borra la línea plantilla.
     También sirve [MATERIALES] o [LISTA MATERIALES] como token suelto.
"""
import copy
import os
from collections import OrderedDict
from io import BytesIO

from pptx import Presentation
from PIL import Image
from googleapiclient.http import MediaIoBaseDownload

from config import IDX
from config_ingresos import (
    NOMBRE_TEMPLATE_INGRESOS,
    NOMBRE_CARPETA_TEMPLATES,
    TEMPLATE_PPT_INGRESOS_ID,
    MAX_FOTOS_INGRESOS,
    FORMATO_LINEA_MATERIAL,
    SUMAR_CANTIDADES_REPETIDAS,
)
from utils import fmt, buscar_foto_blob, parse_fecha
from google_clients import get_drive

# Helpers reutilizados del flujo de Implementación (no se duplica lógica)
from generar_ppt import (
    _normalizar_imagen,
    _reemplazar_texto,
    _encontrar_placeholders_fotos,
    _calcular_posiciones,
    _duplicar_slide,
    _limpiar_placeholder_foto,
)

MIME_SLIDES = "application/vnd.google-apps.presentation"
MIME_PPTX = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
MIME_FOLDER = "application/vnd.google-apps.folder"

NS_A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"

# Cache: el template es el mismo para todas las campañas de la corrida.
_template_cache = None


# ============================================================
# BÚSQUEDA Y DESCARGA DEL TEMPLATE
# ============================================================
def _buscar_template_ingresos(drive):
    """
    Ubica el archivo del template. Orden de intentos:
      1. TEMPLATE_PPT_INGRESOS_ID en config_ingresos.py
      2. Variable de entorno INGRESOS_TEMPLATE_ID
      3. Dentro de la carpeta "Templates_Reportes" (Mi unidad)
      4. Búsqueda global por nombre

    Devuelve (file_id, mime_type).
    """
    id_fijo = (TEMPLATE_PPT_INGRESOS_ID or os.environ.get("INGRESOS_TEMPLATE_ID") or "").strip()
    if id_fijo:
        info = drive.files().get(
            fileId=id_fijo, fields="id, name, mimeType", supportsAllDrives=True
        ).execute()
        print(f"  [TEMPLATE] ID fijo: '{info['name']}' ({info['id']})")
        return info["id"], info["mimeType"]

    nombre_q = NOMBRE_TEMPLATE_INGRESOS.replace("'", "\\'")
    print(f"  [TEMPLATE] Buscando '{NOMBRE_TEMPLATE_INGRESOS}'...")

    # --- Intento 1: dentro de la carpeta Templates_Reportes ---
    try:
        carpeta_q = NOMBRE_CARPETA_TEMPLATES.replace("'", "\\'")
        resp = drive.files().list(
            q=f"name = '{carpeta_q}' and mimeType = '{MIME_FOLDER}' and trashed = false",
            fields="files(id, name)",
            supportsAllDrives=True, includeItemsFromAllDrives=True,
        ).execute()
        for carpeta in resp.get("files", []):
            resp2 = drive.files().list(
                q=(f"name = '{nombre_q}' and '{carpeta['id']}' in parents "
                   f"and trashed = false"),
                fields="files(id, name, mimeType)",
                supportsAllDrives=True, includeItemsFromAllDrives=True,
            ).execute()
            archivos = resp2.get("files", [])
            if archivos:
                a = archivos[0]
                print(f"  [TEMPLATE] ✓ En '{NOMBRE_CARPETA_TEMPLATES}': {a['name']} ({a['id']})")
                return a["id"], a["mimeType"]
        print(f"  [TEMPLATE] No está en '{NOMBRE_CARPETA_TEMPLATES}'. Busco global...")
    except Exception as e:
        print(f"  [TEMPLATE] Error buscando en la carpeta: {e}. Busco global...")

    # --- Intento 2: búsqueda global por nombre ---
    resp3 = drive.files().list(
        q=f"name = '{nombre_q}' and trashed = false",
        fields="files(id, name, mimeType)",
        supportsAllDrives=True, includeItemsFromAllDrives=True,
    ).execute()
    archivos3 = resp3.get("files", [])
    if archivos3:
        a = archivos3[0]
        print(f"  [TEMPLATE] ✓ Encontrado (global): {a['name']} ({a['id']})")
        return a["id"], a["mimeType"]

    raise RuntimeError(
        f"No se encontró '{NOMBRE_TEMPLATE_INGRESOS}'. Revisá que la carpeta "
        f"'{NOMBRE_CARPETA_TEMPLATES}' (o el archivo) esté compartida con la "
        f"cuenta de servicio, o pegá el ID en TEMPLATE_PPT_INGRESOS_ID."
    )


def _descargar_template(drive, file_id, mime):
    """Baja el template como .pptx en memoria (exporta si es Google Slides)."""
    if mime == MIME_SLIDES:
        req = drive.files().export_media(fileId=file_id, mimeType=MIME_PPTX)
    else:
        req = drive.files().get_media(fileId=file_id, supportsAllDrives=True)
    buf = BytesIO()
    downloader = MediaIoBaseDownload(buf, req)
    done = False
    while not done:
        _, done = downloader.next_chunk()
    buf.seek(0)
    return buf


def _template_bytes(drive):
    """Devuelve los bytes del template, cacheados para toda la corrida."""
    global _template_cache
    if _template_cache is None:
        file_id, mime = _buscar_template_ingresos(drive)
        _template_cache = _descargar_template(drive, file_id, mime).getvalue()
        print(f"  [TEMPLATE] Descargado ({len(_template_cache)/1024:.0f} KB)")
    return BytesIO(_template_cache)


# ============================================================
# AGRUPACIÓN POR LOCAL
# ============================================================
def _clave_local(row):
    """
    Clave de agrupación: ID_LOCAL.

    Si viene vacío se cae a "COD|NOMBRE_SALA" para no perder la fila ni
    mezclarla con otro local (todas las filas sin ID caerían en la misma
    bolsa si usáramos una clave fija).
    """
    id_local = str(row[IDX["ID_LOCAL"]] or "").strip()
    if id_local:
        return id_local
    cod = str(row[IDX["COD"]] or "").strip()
    sala = str(row[IDX["NOMBRE_SALA"]] or "").strip()
    return f"SIN_ID|{cod}|{sala}"


def _orden_codigo(row):
    """Orden numérico real de la sala: 'J9' antes que 'J633'."""
    import re
    cod = str(row[IDX["COD"]] or "").strip()
    m = re.match(r"^([A-Za-z]*)(\d+)?", cod)
    if m:
        return (m.group(1) or "", int(m.group(2)) if m.group(2) else 0, cod)
    return (cod, 0, cod)


def agrupar_por_local(filas):
    """
    Agrupa las filas por local. Devuelve OrderedDict {clave: [filas]},
    ordenado por código de sala (numérico real).
    """
    grupos = OrderedDict()
    for row in sorted(filas, key=_orden_codigo):
        grupos.setdefault(_clave_local(row), []).append(row)
    return grupos


def _proceso_limpio(row):
    """Normaliza 'Reagenda interna/externa' → 'Reagenda' (igual que Implementación)."""
    p = str(row[IDX["PROCESO"]] or "").strip()
    return "Reagenda" if p.lower().startswith("reagenda") else p


def materiales_del_local(filas):
    """
    Listado de materiales del local, sin duplicados.

    Devuelve lista de dicts {"material", "cantidad", "proceso"} en el orden en
    que aparecen. Dos filas con el mismo material + cantidad + proceso se
    colapsan (duplicado de carga). Si SUMAR_CANTIDADES_REPETIDAS está en True,
    en vez de colapsar suma las cantidades.
    """
    items = OrderedDict()
    for row in filas:
        material = str(row[IDX["MATERIAL"]] or "").strip()
        cantidad = str(row[IDX["CANTIDAD"]] or "").strip()
        proceso = _proceso_limpio(row)

        # Fila sin material ni proceso: no aporta nada al listado
        if not material and not proceso:
            continue

        clave = (material.lower(), proceso.lower()) if SUMAR_CANTIDADES_REPETIDAS \
            else (material.lower(), cantidad.lower(), proceso.lower())

        if clave in items:
            if SUMAR_CANTIDADES_REPETIDAS:
                previa, nueva = items[clave]["cantidad"], cantidad
                try:
                    items[clave]["cantidad"] = str(int(float(previa or 0)) + int(float(nueva or 0)))
                except ValueError:
                    # Cantidades no numéricas ("2 cajas"): no se pueden sumar
                    items[clave]["cantidad"] = f"{previa} + {nueva}"
            continue

        items[clave] = {"material": material, "cantidad": cantidad, "proceso": proceso}

    return list(items.values())


def fotos_del_local(filas, maximo=None):
    """
    Fotos del local, quitando duplicados.

    Las fotos vienen repetidas en todas las filas del local (una activación
    tiene las mismas 2 fotos cargadas material por material). Se recorren
    todas las filas, se junta todo y se descarta lo repetido comparando el
    NOMBRE DE ARCHIVO en minúsculas: la misma foto puede venir con rutas
    distintas según cómo la haya subido AppSheet.
    """
    if maximo is None:
        maximo = MAX_FOTOS_INGRESOS

    vistas, rutas = set(), []
    for row in filas:
        for k in ("FOTO1", "FOTO2", "FOTO3", "FOTO4"):
            ruta = str(row[IDX[k]] or "").strip()
            if not ruta:
                continue
            huella = ruta.split("/")[-1].strip().lower()
            if huella in vistas:
                continue
            vistas.add(huella)
            rutas.append(ruta)
            if len(rutas) >= maximo:
                return rutas
    return rutas


# ============================================================
# LISTADO DE MATERIALES DENTRO DE LA SLIDE
# ============================================================
def _texto_parrafo(p_el):
    """Texto completo de un párrafo (elemento <a:p>)."""
    return "".join(t.text or "" for t in p_el.iter(f"{NS_A}t"))


def _set_texto_parrafo(p_el, texto):
    """
    Escribe `texto` en el párrafo conservando el formato del primer run.
    Los runs siguientes se vacían (misma técnica que _reemplazar_texto).
    """
    runs = list(p_el.iter(f"{NS_A}r"))
    if not runs:
        return False
    primero = True
    for r in runs:
        for t in r.iter(f"{NS_A}t"):
            t.text = texto if primero else ""
            primero = False
    return True


def _es_linea_materiales(texto):
    """
    ¿Este párrafo es la línea plantilla del listado de materiales?

    Cuenta como tal si tiene [MATERIAL] junto con [CANTIDAD] o [PROCESO],
    o si es el token suelto [MATERIALES] / [LISTA MATERIALES].
    """
    t = texto.upper()
    if "[MATERIALES]" in t or "[LISTA MATERIALES]" in t:
        return True
    return "[MATERIAL]" in t and ("[CANTIDAD]" in t or "[PROCESO]" in t)


def _render_linea(plantilla, item):
    """Reemplaza los tokens de un item de material en la línea plantilla."""
    linea = plantilla
    for token, valor in (
        ("[MATERIAL]", item["material"]),
        ("[CANTIDAD]", item["cantidad"]),
        ("[PROCESO]", item["proceso"]),
    ):
        linea = linea.replace(token, valor)
    return linea


def expandir_listado_materiales(slide, items):
    """
    Convierte la línea plantilla del template en N líneas, una por material.

    Se duplica el párrafo completo (XML), así que se conserva TODO su formato:
    viñeta, fuente, tamaño, color, sangría. Después se borra la plantilla.

    Devuelve la cantidad de líneas escritas, o -1 si no encontró la línea
    plantilla en la slide (se avisa en el log para que se corrija el template).
    """
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            p_el = para._p
            texto = _texto_parrafo(p_el)
            if not _es_linea_materiales(texto):
                continue

            # Si es el token suelto, se usa el formato de config
            plantilla = texto
            if "[MATERIALES]" in texto.upper() or "[LISTA MATERIALES]" in texto.upper():
                plantilla = FORMATO_LINEA_MATERIAL

            if not items:
                p_el.getparent().remove(p_el)
                return 0

            anterior = p_el
            for item in items:
                nuevo = copy.deepcopy(p_el)
                _set_texto_parrafo(nuevo, _render_linea(plantilla, item))
                anterior.addnext(nuevo)
                anterior = nuevo

            p_el.getparent().remove(p_el)   # borrar la línea plantilla
            return len(items)

    return -1


# ============================================================
# TOKENS
# ============================================================
def _tokens_portada(campana, cliente, etiqueta_semana, semana, anio,
                    lunes, domingo, hoy, filas, n_locales, n_materiales):
    periodo = f"{lunes.strftime('%d/%m/%Y')} al {domingo.strftime('%d/%m/%Y')}"
    return [
        ("[FECHA ENVIO DEL REPORTE]", hoy.strftime("%d/%m/%Y")),
        ("[TOTAL MATERIALES]", str(n_materiales)),
        ("[TOTAL REGISTROS]", str(len(filas))),
        ("[TOTAL LOCALES]", str(n_locales)),
        ("[TOTAL SALAS]", str(n_locales)),
        ("[N SEMANA]", str(semana)),
        ("[CAMPAÑA]", campana),
        ("[CAMPANA]", campana),
        ("[CLIENTE]", cliente),
        ("[PERIODO]", periodo),
        ("[SEMANA]", etiqueta_semana),
        ("[AÑO]", str(anio)),
        ("[ANIO]", str(anio)),
        ("[FECHA]", hoy.strftime("%d/%m/%Y")),
    ]


def _rango_fechas(filas, campo):
    """
    Fecha representativa del local para un campo dado.
    FECHA_INICIO → la mínima; el resto → la máxima. Las filas del mismo local
    pueden traer fechas distintas por material, y el local se activa en un
    rango, no en un día.
    """
    fechas = [parse_fecha(r[IDX[campo]]) for r in filas]
    fechas = [f for f in fechas if f]
    if not fechas:
        return ""
    elegida = min(fechas) if campo == "FECHA_INICIO" else max(fechas)
    return elegida.strftime("%d/%m/%Y")


def _primero_no_vacio(filas, campo):
    for r in filas:
        v = str(r[IDX[campo]] or "").strip()
        if v:
            return v
    return ""


def _unicos(filas, campo, transformar=None):
    """Valores únicos de una columna entre las filas del local, en orden."""
    vistos, out = set(), []
    for r in filas:
        v = transformar(r) if transformar else str(r[IDX[campo]] or "").strip()
        if v and v.lower() not in vistos:
            vistos.add(v.lower())
            out.append(v)
    return out


def _tokens_local(filas, etiqueta_semana, items):
    """
    Tokens de la slide de un local.

    OJO CON EL ORDEN: los combinados van PRIMERO. Si reemplazáramos [CÓD.]
    antes que "[CÓD.] - [NOMBRE SALA]", el combinado ya no existiría.

    [MATERIAL] / [CANTIDAD] / [PROCESO] se resuelven como RESUMEN del local
    (únicos concatenados y total). El listado línea por línea ya se expandió
    antes con expandir_listado_materiales(), así que acá solo se cubren
    apariciones sueltas de esos tokens en otra parte de la slide.
    """
    base = filas[0]
    g = lambda k: str(base[IDX[k]] or "")  # noqa: E731
    cod, sala = g("COD"), g("NOMBRE_SALA")
    direccion, comuna, region = g("DIRECCION"), g("COMUNA"), g("REGION")

    materiales = ", ".join(i["material"] for i in items if i["material"])
    procesos = ", ".join(_unicos(filas, None, transformar=_proceso_limpio))
    try:
        total_cant = str(int(sum(float(i["cantidad"] or 0) for i in items)))
    except ValueError:
        total_cant = ""

    return [
        # --- combinados (primero) ---
        ("[CÓD.] - [NOMBRE SALA]", f"{cod} - {sala}"),
        ("[COD] - [NOMBRE SALA]", f"{cod} - {sala}"),
        ("[DIRECCIÓN], [COMUNA] - [REGION]", f"{direccion}, {comuna} - {region}"),
        ("[DIRECCION], [COMUNA] - [REGION]", f"{direccion}, {comuna} - {region}"),
        # --- individuales ---
        ("[TIPO ACTIVIDAD]", g("TIPO_ACTIVIDAD")),
        ("[TOTAL MATERIALES]", str(len(items))),
        ("[NOMBRE SALA]", sala),
        ("[FECHA COMPROMISO]", _rango_fechas(filas, "FECHA_COMPROMISO")),
        ("[FECHA TERMINO]", _rango_fechas(filas, "FECHA_TERMINO")),
        ("[FECHA ENTREGA]", _rango_fechas(filas, "FECHA_ENTREGA")),
        ("[FECHA INICIO]", _rango_fechas(filas, "FECHA_INICIO")),
        ("[HORA TERMINO]", _primero_no_vacio(filas, "HORA_TERMINO")),
        ("[HORA INICIO]", _primero_no_vacio(filas, "HORA_INICIO")),
        ("[OBSERVACIONES]", _primero_no_vacio(filas, "OBSERVACIONES")),
        ("[ID LOCAL]", g("ID_LOCAL")),
        ("[DIRECCIÓN]", direccion),
        ("[DIRECCION]", direccion),
        ("[CATEGORIA]", ", ".join(_unicos(filas, "CATEGORIA"))),
        ("[CATEGORÍA]", ", ".join(_unicos(filas, "CATEGORIA"))),
        ("[EJECUTIVO]", g("EJECUTIVO")),
        ("[ACTIVIDAD]", g("ACTIVIDAD")),
        ("[CANTIDAD]", total_cant),
        ("[MATERIAL]", materiales),
        ("[PROCESO]", procesos),
        ("[DETALLE]", _primero_no_vacio(filas, "DETALLE")),
        ("[CADENA]", g("CADENA")),
        ("[COMUNA]", comuna),
        ("[REGION]", region),
        ("[REGIÓN]", region),
        ("[SEMANA]", etiqueta_semana),
        ("[MARCA]", ", ".join(_unicos(filas, "MARCA"))),
        ("[GUIA]", ", ".join(_unicos(filas, "GUIA"))),
        ("[GUÍA]", ", ".join(_unicos(filas, "GUIA"))),
        ("[CÓD.]", cod),
        ("[COD]", cod),
        ("[SKU]", ", ".join(_unicos(filas, "SKU"))),
    ]


# ============================================================
# FOTOS EN LA SLIDE
# ============================================================
def _posiciones_fotos(n_fotos, placeholders_info):
    """
    Si hay exactamente un placeholder por foto, se usan las posiciones tal
    cual las dejó el template (respeta separación y tamaños del diseño).
    Si no coinciden, se cae al cálculo de grilla de Implementación.
    """
    if n_fotos == len(placeholders_info):
        return placeholders_info
    return _calcular_posiciones(n_fotos, placeholders_info)


def _insertar_fotos(slide, drive, rutas):
    """Inserta las fotos en los placeholders, manteniendo proporción."""
    placeholders = _encontrar_placeholders_fotos(slide)
    info = [(sh.left, sh.top, sh.width, sh.height)
            for _, sh in sorted(placeholders.items())]
    for sh in placeholders.values():
        _limpiar_placeholder_foto(sh)

    if not rutas or not info:
        return 0

    posiciones = _posiciones_fotos(len(rutas), info)
    puestas = 0
    for i, ruta in enumerate(rutas):
        if i >= len(posiciones):
            break
        area_left, area_top, area_w, area_h = posiciones[i]
        blob = buscar_foto_blob(drive, ruta)
        if blob is None:
            print(f"      [WARN] Foto {i+1} no encontrada: {ruta}")
            continue
        try:
            norm = _normalizar_imagen(blob)
            if norm is None:
                print(f"      [SKIP] Foto {i+1} omitida")
                continue
            norm.seek(0)
            img_w, img_h = Image.open(norm).size
            norm.seek(0)

            # "Fit" dentro del área, manteniendo proporción y centrado
            if (img_w / img_h) > (area_w / area_h):
                final_w, final_h = area_w, int(area_w / (img_w / img_h))
            else:
                final_h, final_w = area_h, int(area_h * (img_w / img_h))

            slide.shapes.add_picture(
                norm,
                area_left + (area_w - final_w) // 2,
                area_top + (area_h - final_h) // 2,
                width=final_w, height=final_h,
            )
            puestas += 1
            print(f"      OK foto {i+1}")
        except Exception as e:
            print(f"      [ERROR] foto {i+1}: {e}")
    return puestas


# ============================================================
# GENERACIÓN
# ============================================================
def generar_ppt_ingresos(campana, filas, etiqueta_semana, semana, anio,
                         lunes, domingo, hoy):
    """
    Genera el PPT semanal de ingresos de una campaña: UNA SLIDE POR LOCAL.
    Devuelve los bytes del .pptx.
    """
    drive = get_drive()
    cliente = str(filas[0][IDX["CLIENTE"]] or "")
    pres = Presentation(_template_bytes(drive))

    if len(pres.slides) < 2:
        raise RuntimeError(
            f"El template '{NOMBRE_TEMPLATE_INGRESOS}' tiene "
            f"{len(pres.slides)} slide(s). Se necesitan 2: portada + local."
        )

    # ---- AGRUPAR POR LOCAL ----
    locales = agrupar_por_local(filas)
    items_por_local = {k: materiales_del_local(v) for k, v in locales.items()}
    total_materiales = sum(len(v) for v in items_por_local.values())
    print(f"  Locales: {len(locales)} | materiales: {total_materiales} "
          f"| filas: {len(filas)}")

    # ---- PORTADA ----
    portada = pres.slides[0]
    for buscar, reemplazar in _tokens_portada(
        campana, cliente, etiqueta_semana, semana, anio, lunes, domingo, hoy,
        filas, len(locales), total_materiales
    ):
        _reemplazar_texto(portada, buscar, reemplazar)

    # ---- UNA SLIDE POR LOCAL ----
    template_local = pres.slides[1]
    slides = [template_local]
    for _ in range(len(locales) - 1):
        slides.append(_duplicar_slide(pres, template_local))

    aviso_listado = False
    for i, (clave, filas_local) in enumerate(locales.items()):
        slide = slides[i]
        items = items_por_local[clave]
        base = filas_local[0]
        cod = str(base[IDX["COD"]] or "")
        sala = str(base[IDX["NOMBRE_SALA"]] or "")
        print(f"    Slide {i+1}/{len(locales)}: {cod} {sala} "
              f"({len(filas_local)} filas → {len(items)} materiales)")

        # 1) El listado PRIMERO: si no, el reemplazo general pisaría los
        #    tokens de la línea plantilla con el resumen del local.
        escritas = expandir_listado_materiales(slide, items)
        if escritas == -1:
            aviso_listado = True
        else:
            for it in items:
                print(f"        - {it['material']} ({it['cantidad']}) : {it['proceso']}")

        # 2) Tokens del local
        for buscar, reemplazar in _tokens_local(filas_local, etiqueta_semana, items):
            _reemplazar_texto(slide, buscar, reemplazar)

        # 3) Fotos, sin duplicados entre las filas del local
        rutas = fotos_del_local(filas_local)
        if len(filas_local) > 1:
            print(f"        Fotos: {len(rutas)} única(s) de {len(filas_local)} filas")
        _insertar_fotos(slide, drive, rutas)

    if aviso_listado:
        print("  [AVISO] No se encontró la línea del listado de materiales en el "
              "template. Agregá una línea con '[MATERIAL] ([CANTIDAD]) : [PROCESO]' "
              "en la slide 2.")

    buf = BytesIO()
    pres.save(buf)
    return buf.getvalue()
